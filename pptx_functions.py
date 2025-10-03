from pptx.dml.color import RGBColor
import random
from pptx.oxml.xmlchemy import OxmlElement
from pptx import Presentation
import fitz  # PyMuPDF
import comtypes.client
from comtypes import CoInitialize, CoUninitialize
import os
from PIL import Image as PILImage
from pptx.dml.color import RGBColor
import re
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt, Inches, Cm
from datetime import datetime, timedelta, date
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE, MSO_CONNECTOR

colors = {
    "green_text": RGBColor(112, 173, 71),
    "red_text": RGBColor(239, 107, 107),
    "red_back": RGBColor(248, 203, 173),
    "red_back_map": RGBColor(248, 203, 173),
}


def df_to_chart_data(df, delete=False):
    if delete:
        df = df.drop(df.columns[1], axis=1)
    if df.shape[1] != 2:
        raise ValueError("DataFrame должен содержать ровно 2 столбца: (категории, значения)")

    return list(df.itertuples(index=False, name=None))


def format_category_text(category):
    """Извлекает текст между кавычками «...» и форматирует его:
    - Первые буквы всех слов заглавные
    - 'район' и 'района' остаются в нижнем регистре
    """

    # Извлекаем текст между кавычками «...»
    match = re.search(r"«([^»]+)»", category)
    if not match:
        return None  # Если кавычек нет, возвращаем None

    text = match.group(1).lower()  # Приводим весь текст к нижнему регистру

    # Разбиваем на слова, приводим каждое к Title Case, кроме "район" и "района"
    words = text.split()
    formatted_words = [word.capitalize() if word not in ["район", "района"] else word for word in words]

    return " ".join(formatted_words)


def update_diagramms(input_path, output_path, diagramms_data):
    print(diagramms_data, "aaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaaa")
    comtypes.CoInitialize()  # Инициализация COM (важно!)
    ppt = comtypes.client.CreateObject("PowerPoint.Application")
    ppt.Visible = True  # Можно скрыть, установив False
    # Открываем презентацию
    ppt_path = os.path.abspath(input_path)
    output_path = os.path.abspath(output_path)
    presentation = ppt.Presentations.Open(ppt_path, WithWindow=False)
    presentation = region_data_update(presentation, 2, diagramms_data["theme0_earliest"],
                                      diagramms_data["theme0_latest"], 1)
    presentation = region_data_update(presentation, 3, diagramms_data["theme1_earliest"],
                                      diagramms_data["theme1_latest"], 4)
    presentation = region_data_update(presentation, 4, diagramms_data["theme2_earliest"],
                                      diagramms_data["theme2_latest"], 4)
    for i in range(3):
        shape_dates = 58 if i == 0 else 44
        presentation = date_diagramm_update(presentation, i + 2, diagramms_data[f"theme{i}_dates"],
                                            diagramms_data[f"theme{i}_values"], shape_dates)
        for o in range(3):
            if i == 0:
                shape_id = 28 if o == 0 else 32 if o == 1 else 36
            else:
                shape_id = 15 if o == 0 else 19 if o == 1 else 23
            presentation = org_diagramm_update(presentation, i + 2, diagramms_data[f"theme{i}_okrug{o + 1}_resp"],
                                               diagramms_data[f"theme{i}_okrug{o + 1}_early"],
                                               diagramms_data[f"theme{i}_okrug{o + 1}_late"],
                                               shape_id)

    presentation.SaveAs(output_path)
    presentation.Close()
    ppt.Quit()
    comtypes.CoUninitialize()


#     1 -- regions
#     31
#     35
#     39
#     61 -- dates


def excel_date(dt):
    """Преобразует дату в формат Excel (количество дней с 30.12.1899)"""
    base_date = datetime(1899, 12, 30).date()  # Приводим к date
    return (dt - base_date).days


def org_diagramm_update(presentation, slide_id, categories, previous_values, current_values, shape_id):
    # Выбираем слайд
    slide = presentation.Slides(slide_id + 1)  # COM использует 1-based индексацию

    # Ищем диаграмму на слайде
    for sh, shape in enumerate(slide.Shapes):
        if shape.HasChart:
            print("ShapeID : ", sh)
        if shape.HasChart and sh == shape_id:
            chart = shape.Chart
            workbook = chart.ChartData.Workbook  # Открываем данные диаграммы
            sheet = workbook.Worksheets(1)  # Первый лист с данными

            # Убираем защиту, если есть
            try:
                sheet.Unprotect()
            except:
                pass  # Если нет защиты, просто продолжаем
            sheet.Cells.ClearContents()
            # Заполняем заголовки
            sheet.Cells(1, 1).FormulaR1C1 = "Организация"
            sheet.Cells(1, 2).FormulaR1C1 = "Предыдущий период"
            sheet.Cells(1, 3).FormulaR1C1 = "Отчетный период"

            # Вставляем новые данные
            for row_idx, (category, prev_val, curr_val) in enumerate(zip(categories, previous_values, current_values),
                                                                     start=2):
                sheet.Cells(row_idx, 1).FormulaR1C1 = format_category_text(category)
                sheet.Cells(row_idx, 2).FormulaR1C1 = prev_val
                sheet.Cells(row_idx, 3).FormulaR1C1 = curr_val

            # Обновляем диаграмму
            chart.ChartData.Activate()
            workbook.Close(SaveChanges=True)  # Закрываем Excel-данные
            chart.Refresh()
    return presentation


def date_diagramm_update(presentation, slide_id, dates, values, shape_id):
    # Выбираем слайд
    slide = presentation.Slides(slide_id + 1)  # COM использует 1-based индексацию
    # Ищем диаграмму на слайде
    for sh, shape in enumerate(slide.Shapes):
        if shape.HasChart and sh == shape_id:
            chart = shape.Chart
            workbook = chart.ChartData.Workbook
            sheet = workbook.Worksheets(1)

            # Отключаем защиту листа (если есть)
            try:
                sheet.Unprotect()
            except:
                pass  # Если нет защиты, просто продолжаем

            # Обновляем данные
            data = [["Дата", "Значение"]] + list(zip(dates, values))

            for row_idx, row in enumerate(data, start=1):
                for col_idx, value in enumerate(row, start=1):
                    if isinstance(value, datetime) or isinstance(value, (datetime, date)):
                        sheet.Cells(row_idx, col_idx).FormulaR1C1 = excel_date(value)  # Даты в Excel-формате
                    else:
                        sheet.Cells(row_idx, col_idx).FormulaR1C1 = value  # Обычные значения

            max_value = max(values)
            max_index = values.index(max_value)

            # **Изменяем ось X (горизонтальную ось)**
            axis_x = chart.Axes(1)  # xlCategory = 1
            axis_x.MinimumScale = excel_date(dates[0])  # Минимальная дата
            axis_x.MaximumScale = excel_date(dates[-1])  # Максимальная дата
            # axis_x.TickLabels.NumberFormat = "DD.MM.YYYY"  # Формат подписей оси
            chart.ChartData.Activate()
            workbook.Close(SaveChanges=True)
            chart.Refresh()

            # Найти одно максимальных значения и их индексы

            # Определить координаты точки на диаграмме
            chart_area = shape.Chart.PlotArea
            # 823 874
            x_position = 618 + 41 * max_index
            y_position = chart_area.Top + chart_area.Height + 133  # Немного ниже диаграммы

            # Добавление линии под максимальным значением
            line = slide.Shapes.AddLine(
                x_position, y_position, x_position + 35, y_position
            )
            line.Line.ForeColor.RGB = 0  # Черный цвет линии
            line.Line.Weight = 3  # Толщина линии

            # Сообщение в консоль
            print(f"Диаграмма обновлена: {shape.Name}, Дата макс. значения: {dates[max_index]}, Значение: {max_value}")

    return presentation


def region_data_update(presentation, slide_id, previous_period, current_period, shape_id):
    categories = ["ЦАО", "САО", "СВАО", "ВАО", "ЮВАО", "ЮАО", "ЮЗАО", "ЗАО", "СЗАО", "ЗелАО", "ТиНАО", "ГБУ«АВД»", "Иные"]
    slide = presentation.Slides(slide_id + 1)  # COM использует 1-based индексацию
    for sh, shape in enumerate(slide.Shapes):
        if shape.HasChart and sh == shape_id:  # Проверяем, является ли объект диаграммой
            chart = shape.Chart
            workbook = chart.ChartData.Workbook  # Открываем таблицу Excel внутри диаграммы
            sheet = workbook.Worksheets(1)  # Первый лист с данными диаграммы
            try:
                sheet.Unprotect()
            except:
                pass  # Если нет защиты, просто продолжаем

            # Обновляем данные
            data = [["Район", "Предыдущий период", "Отчетный период"]] + list(
                zip(categories, previous_period, current_period))

            for row_idx, row in enumerate(data, start=1):
                for col_idx, value in enumerate(row, start=1):
                    sheet.Cells(row_idx, col_idx).FormulaR1C1 = value  # Используем FormulaR1C1
            chart.ChartData.Activate()
            workbook.Close(SaveChanges=True)
            chart.Refresh()

            # Найти три максимальных значения и их индексы
            top_values = sorted(zip(current_period, categories), reverse=True)[:3]
            top_indices = [categories.index(cat) for _, cat in top_values]

            # Определить координаты колонок диаграммы
            axis = chart.Axes(1)  # X-ось
            chart_area = shape.Chart.PlotArea

            for idx in top_indices:
                # Вычисление координат (грубая оценка, зависит от типа диаграммы)
                x_position = chart_area.Left + (idx + 1) * (chart_area.Width / len(categories))
                y_position = chart_area.Top + chart_area.Height + 257  # Немного ниже диаграммы

                # Добавление линии
                line = slide.Shapes.AddLine(
                    x_position - 10, y_position, x_position + 50, y_position
                )
                line.Line.ForeColor.RGB = 0  # Черный цвет линии
                line.Line.Weight = 2.7  # Толщина линии

    return presentation


# Деинициализация COM


def get_top_rows_with_ties(counts_list, top_n=3):
    sorted_list = sorted(counts_list, key=lambda x: float(f"{x[1]}".replace(" ", "")), reverse=True)
    top_with_ties = []
    current_rank = 0
    last_value = None

    for row in sorted_list:
        if current_rank < top_n or (row[1] == last_value):
            top_with_ties.append(row[0])
            current_rank += 1
            if row[1] != last_value:
                last_value = row[1]
        else:
            break
    if len(sorted_list) > top_n:
        third_value = sorted_list[top_n - 1][1]
        fourth_value = sorted_list[top_n][1]
        if third_value != fourth_value:
            top_with_ties = [row[0] for row in sorted_list[:top_n]]

    return top_with_ties


def find_in_group(presentation):
    for sl, slide in enumerate(presentation.slides):
        for sh, shape in enumerate(slide.shapes):
            if shape.shape_type == 6:  # Группа
                for gr, grouped_shape in enumerate(shape.shapes):
                    if grouped_shape.has_text_frame:
                        for p, paragraph in enumerate(grouped_shape.text_frame.paragraphs):
                            text = ""
                            for i, run in enumerate(paragraph.runs):
                                text = f"{text}{run.text}"
                            # print(f"sh: {sh} gr: {gr} p: {p} : text = {text}")


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def _set_shape_transparency(shape, alpha):
    """ Set the transparency (alpha) of a shape"""
    ts = shape.fill._xPr.solidFill
    sF = ts.get_or_change_to_srgbClr()
    sE = SubElement(sF, 'a:alpha', val=str(alpha))


def color_shapes_on_slide(presentation, slide_index):
    slide = presentation.slides[slide_index]  # Получаем указанный слайд (в данном случае 2-й, индекс 1)
    shape_colors = []  # Список для хранения номеров фигур и их цветов

    for sh, shape in enumerate(slide.shapes):
        # Проверяем, поддерживает ли фигура изменение заливки
        if not hasattr(shape, "fill") or shape.fill is None:
            continue  # Пропускаем фигуры без атрибута `fill`

        # Генерируем случайный цвет
        random_color = RGBColor(random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))

        # Применяем цвет к фигуре
        shape.fill.solid()
        shape.fill.fore_color.rgb = random_color
        shape_colors.append((sh, f"#{random_color[0]:02x}{random_color[1]:02x}{random_color[2]:02x}"))

    # Печатаем массив с номерами и цветами
    print(shape_colors)
    return presentation


def hard_code_color(presentation, type_report):
    # Открытие презентации
    if type_report != "week":
        s1 = 0
        s2 = 26
    else:
        s1 = 0
        s2 = 18
    #     Красит текст а должно контур на первом слайде
    for sl, slide in enumerate(presentation.slides):
        for s, shape in enumerate(slide.shapes):
            # Проверка, является ли shape текстовым объектом
            if shape.has_text_frame:
                for i, paragraph in enumerate(shape.text_frame.paragraphs):
                    # for r, run in enumerate(paragraph.runs):
                    #     print("slide", sl, "shape", s, " Parag", i, "run", r, "text = ", run.text)
                    if (sl == s1) and (s == s2):
                        color = colors["red_text"]
                        for run in paragraph.runs:

                            if "-" in run.text:
                                color = colors["green_text"]
                        for run in paragraph.runs:
                            run.font.color.rgb = color
                    if type_report == "week":
                        if sl in [2, 3, 4, 5] and s in [23, 53]:
                            color = colors["red_text"]
                            for run in paragraph.runs:

                                if "-" in run.text:
                                    color = colors["green_text"]
                            for r, run in enumerate(paragraph.runs):
                                if r == 0:
                                    run.font.color.rgb = RGBColor(255, 255, 255)
                                else:
                                    run.font.color.rgb = color
    return presentation


def find_text_box(pptx_file_path):
    # Открытие презентации
    presentation = Presentation(pptx_file_path)
    for sl, slide in enumerate(presentation.slides):
        for s, shape in enumerate(slide.shapes):
            # Проверка, является ли shape текстовым объектом
            if shape.has_text_frame:
                for i, paragraph in enumerate(shape.text_frame.paragraphs):
                    print("slide", sl, "shape", s, " Parag", i)
                    for run in paragraph.runs:
                        print(run.text)


def find_text_table(presentation, type_report):
    counts_list_col2 = []  # Список для хранения строк и чисел из 2-й колонки
    counts_list_col3 = []  # Список для хранения строк и чисел из 3-й колонки
    counts_list_col4 = []  # Список для хранения строк и чисел из 4-й колонки
    if type_report != "week":
        shape_to_r_map = {
            34: 1,
            17: 2,
            14: 3,
            11: 4,
            46: 5,
            42: 6,
            38: 7,
            5: 8,
            8: 9,
            24: 10,
            30: 11,
        }
        shape_to_r_map_text = {
            35: 1,
            18: 2,
            15: 3,
            12: 4,
            47: 5,
            43: 6,
            39: 7,
            6: 8,
            9: 9,
            25: 10,
            33: 11,
        }
        sh_by_type = 51
        slide_for_all_table = 2
    else:
        shape_to_r_map = {
            41: 1,
            26: 2,
            23: 3,
            20: 4,
            53: 5,
            49: 6,
            45: 7,
            14: 8,
            17: 9,
            33: 10,
            37: 11,
        }

        shape_to_r_map_text = {
            42: 1,  # ЦАО
            27: 2,  # САО
            24: 3,  # СВАО
            21: 4,  # ВАО
            54: 5,  # ЮВАО
            50: 6,  # ЮАО
            46: 7,  # ЮЗАО
            15: 8,  # ЗАО
            18: 9,  # СЗАО
            34: 10,  # ЗелАО
            38: 11,  # ТиНАО
        }
        sh_by_type = 11
        slide_for_all_table = 5
    for sl, slide in enumerate(presentation.slides):
        if sl == 1:
            for sh, shape in enumerate(slide.shapes):
                if shape.shape_type == 19:  # Таблица
                    table = shape.table
                    last_row_index = len(table.rows) - 1  # Индекс последнего ряда

                    for r, row in enumerate(table.rows):
                        if r == last_row_index:  # Пропустить последний ряд
                            continue

                        for c, cell in enumerate(row.cells):
                            if cell.text_frame:
                                for p, paragraph in enumerate(cell.text_frame.paragraphs):
                                    text = ""
                                    for i, run in enumerate(paragraph.runs):
                                        text = f"{text}{run.text}"
                                    # print(f"sh: {sh} r: {r} c: {c}, p: {p} : text = {text}")
                                    # Условие для 2-й колонки
                                    if c == 2 and sh == sh_by_type:
                                        text = text.replace(" ", "")
                                        numbers = re.findall(r'\b\d+\b', text)
                                        if numbers:
                                            max_number = max(map(int, numbers))
                                            counts_list_col2.append((r, max_number))

                                    # Условие для 3-й колонки
                                    if c == 3 and sh == sh_by_type:
                                        text = text.replace(" ", "").replace(",", "")
                                        numbers = re.findall(r'\b\d+\b', text)
                                        if numbers:
                                            max_number = max(map(int, numbers))
                                            counts_list_col3.append((r, max_number))

                                    # Условие для 4-й колонки
                                    if c == 4 and sh == sh_by_type:
                                        text = text.replace(" ", "").replace(",", "")
                                        numbers = re.findall(r'\b\d+\b', text)
                                        if numbers:
                                            max_number = max(map(int, numbers))
                                            counts_list_col4.append((r, max_number))

                    # Сортировка и выбор топ-3 строк для каждой колонки
                    top_3_rows_col2 = get_top_rows_with_ties(counts_list_col2, top_n=3)
                    top_3_rows_col3 = get_top_rows_with_ties(counts_list_col3, top_n=3)
                    top_3_rows_col4 = get_top_rows_with_ties(counts_list_col4, top_n=3)

                    # Изменение цвета текста для топ-3 строк в 2-й и 3-й колонках
                    for r, row in enumerate(table.rows):
                        if r == last_row_index:  # Пропустить последний ряд
                            continue

                        for c, cell in enumerate(row.cells):
                            if cell.text_frame:
                                for p, paragraph in enumerate(cell.text_frame.paragraphs):
                                    # Изменение для 2-й колонки
                                    if c == 2 and r in top_3_rows_col2:
                                        for i, run in enumerate(paragraph.runs):
                                            run.font.color.rgb = colors["red_text"]  # Красный цвет текста

                                    # Изменение для 3-й колонки
                                    if c == 3 and r in top_3_rows_col3:
                                        for i, run in enumerate(paragraph.runs):
                                            run.font.color.rgb = colors["red_text"]  # Синий цвет текста

                            # Изменение заливки для 4-й колонки
                            if c == 4 and r in top_3_rows_col4:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = colors["red_back"]  # Зеленая заливка ячейки
            for sh, shape in enumerate(slide.shapes):
                if sl == 1:
                    # Проверяем, поддерживает ли фигура изменение заливки
                    if not hasattr(shape, "fill") or shape.fill is None:
                        continue  # Пропускаем фигуры без атрибута `fill`
                    if sh in shape_to_r_map:
                        r = shape_to_r_map[sh]  # Получаем значение `r` для фигуры

                        # Проверяем, находится ли `r` в топ-3 строк
                        if r in top_3_rows_col4:
                            # Если да, изменяем цвет фигуры
                            fill = shape.fill
                            fill.solid()  # Устанавливаем градиентную заливку

                            fill.fore_color.rgb = colors["red_back_map"]
                            _set_shape_transparency(shape, 44000)
            for sh, shape in enumerate(slide.shapes):
                if shape.shape_type == 19:  # Таблица
                    if sh in shape_to_r_map_text:
                        table = shape.table
                        row = table.rows[1]
                        r = shape_to_r_map_text[sh]
                        if r in top_3_rows_col3:
                            cell = row.cells[0]
                            if cell.text_frame:
                                for p, paragraph in enumerate(cell.text_frame.paragraphs):
                                    text = ""
                                    for i, run in enumerate(paragraph.runs):
                                        text = f"{text}{run.text}"
                                        run.font.color.rgb = colors["red_text"]
                                    # print(f"sh: {sh} r: {r} , p: {p} : text = {text}")

        if sl == slide_for_all_table:  # Индекс слайда 2 (третий слайд, 0-индексация)
            for sh, shape in enumerate(slide.shapes):
                if shape.shape_type == 19:  # Проверяем, что это таблица
                    table = shape.table

                    last_row_index = len(table.rows) - 1  # Индекс последней строки
                    last_col_index = len(table.columns) - 1  # Индекс последнего столбца

                    for r, row in enumerate(table.rows):
                        if r == last_row_index or r == 0:  # Пропускаем последнюю строку
                            continue

                        # Создаем список (столбец, значение) для текущей строки
                        counts_list_row = []
                        for c, cell in enumerate(row.cells):
                            if c == last_col_index or c == 0:  # Пропускаем последнюю колонку
                                continue

                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    text = "".join(run.text for run in paragraph.runs)
                                    # Извлекаем числа из текста
                                    text = text.replace(" ", "")
                                    numbers = re.findall(r'\b\d+\b', text)
                                    if numbers:
                                        number = numbers[0]
                                        counts_list_row.append((c, number))  # Сохраняем индекс столбца и значение
                        print(counts_list_row)
                        # Используем вашу функцию для определения топ-3 столбцов
                        top_columns = get_top_rows_with_ties(counts_list_row, top_n=3)
                        print(f"topcollums: {top_columns}")
                        # Окрашиваем ячейки с топовыми значениями
                        for c in top_columns:
                            cell = row.cells[c]
                            if cell.text_frame:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.color.rgb = colors["red_text"]
    return presentation


def color_arrows_in_shape(shape):
    # Обход всех слайдов в презентации

    # Проверка, является ли shape текстовым объектом
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.text == "↗":
                    run.font.color.rgb = colors["red_text"]
                if run.text == "↘":
                    run.font.color.rgb = colors["green_text"]

    # Проверка, является ли shape таблицей
    if shape.shape_type == 19:  # Таблица
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text == "↗":
                                run.font.color.rgb = colors["red_text"]
                            if run.text == "↘":
                                run.font.color.rgb = colors["green_text"]
    # Проверка, является ли shape группой объектов
    if shape.shape_type == 6:  # Группа
        for grouped_shape in shape.shapes:
            if grouped_shape.has_text_frame:
                for paragraph in grouped_shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text == "↗":
                            run.font.color.rgb = colors["red_text"]
                        if run.text == "↘":
                            run.font.color.rgb = colors["green_text"]
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            color_arrows_in_shape(sub_shape)


def runs_from_pptx(pptx_file_path, top_dict, type_report):
    # Открытие презентации
    presentation = Presentation(pptx_file_path)
    presentation = hard_code_color(presentation, type_report)
    presentation = find_text_table(presentation, type_report)
    if top_dict:
        presentation = print_tops_on_second_slide(presentation, top_dict)
    for slide in presentation.slides:
        for shape in slide.shapes:
            color_arrows_in_shape(shape)

    # presentation = color_shapes_on_slide(presentation, 2)
    # presentation = color_shapes_on_slide(presentation, 3)
    # presentation = color_shapes_on_slide(presentation, 4)
    return presentation


def remove(table, column):
    col_idx = table._tbl.tblGrid.index(column._gridCol)
    for tr in table._tbl.tr_lst:
        tr.remove(tr.tc_lst[col_idx])
    table._tbl.tblGrid.remove(column._gridCol)


def remove_columns(pres):
    for slide in pres.slides:
        for shp in slide.shapes:
            if shp.has_table:
                table = shp.table

                # Ширина (сумма ширин) всех столбцов до удаления
                width_before = sum(col.width for col in table.columns)

                # Определяем, какие колонки нужно удалить
                num_cols = len(table.columns)
                cols_to_remove = set()

                # Поиск колонок, содержащих "@"
                for col_idx in range(num_cols):
                    for row in table.rows:
                        cell = row.cells[col_idx]
                        if cell.text_frame:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if "@" in run.text:
                                        cols_to_remove.add(col_idx)

                # Удаление колонок в обратном порядке
                for col_idx in sorted(cols_to_remove, reverse=True):
                    remove(table, table.columns[col_idx])

                # Ширина всех столбцов после удаления
                width_after = sum(col.width for col in table.columns)

                # Разница ширин
                diff = width_before - width_after

                # Сдвигаем таблицу вправо на половину разницы
                # (если diff > 0, значит таблица стала уже, и мы смещаем её вправо)
                shp.left = int(shp.left + diff / 2)

    return pres


def remove_slides_tinao(file_path, tinao_len):
    pres = Presentation(file_path)
    if tinao_len < 6:
        tinao_len = 6
    if tinao_len > 14:
        tinao_len = 14

    slides_to_remove = [11, 12, 13, 14, 15, 16, 17, 18, 19, 20, 21, 22]
    slides_to_remove.remove(11 + tinao_len - 6)
    # Удаляем слайды в обратном порядке, чтобы не нарушить индексы
    for slide_index in sorted(slides_to_remove, reverse=True):
        xml_slides = pres.slides._sldIdLst
        xml_slides.remove(xml_slides[slide_index])
    pres = remove_columns(pres)
    return pres


def convert_pptx_to_pdf(input_file, output_file):
    try:
        CoInitialize()
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.WindowState = 2  # 2 - минимизация окна PowerPoint
        powerpoint.Visible = 1

        presentation = powerpoint.Presentations.Open(os.path.abspath(input_file), WithWindow=False)
        presentation.SaveAs(os.path.abspath(output_file), 32)  # 32 - это формат PDF
        presentation.Close()

    except Exception as e:
        print(f"Ошибка: {e}")
        return False
    finally:
        if 'powerpoint' in locals():
            powerpoint.Quit()
        CoUninitialize()  # Освобождение COM
    return True


def pdf_to_png(pdf_path, output_folder, dpi=150, scale=0.75):
    """Конвертирует PDF в PNG, уменьшает изображение на 25%"""
    os.makedirs(output_folder, exist_ok=True)
    doc = fitz.open(pdf_path)
    slide_paths = []

    for i, page in enumerate(doc):
        output_path = os.path.join(output_folder, f"slide_{i + 1}.png")
        pix = page.get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72))
        img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples)

        # Уменьшаем изображение
        new_size = (int(img.width * scale), int(img.height * scale))
        img = img.resize(new_size)
        img.save(output_path)

        slide_paths.append(output_path)
        print(f"Сохранён слайд: {output_path}")

    return slide_paths


def print_tops_on_second_slide(prs, tops):
    """
    prs  - объект презентации (pptx.Presentation),
    tops - словарь вида:
        {
          "ao1": [("ВАО", 123), ("СВАО", 25), ("ЦАО", 13)],
          "ao2": [("ВАО", 150), ("ЦАО", 85)],
          "ao3": [...],
        }
    """

    # Получаем второй слайд (index=1) — предполагается, что он уже существует
    slide = prs.slides[1]

    # Если нужно обновить заголовок второго слайда
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Данные по округам"

    left_group_box = Cm(14.76)  # Левый отступ для заголовка группы
    top_start = Cm(24.3)  # Начальный верхний отступ (первая строка)
    group_box_width = Cm(0)  # Ширина TextBox под название группы
    group_box_height = Cm(1)

    item_box_height = Cm(1.22)
    horizontal_gap = Cm(0)  # Горизонтальный зазор между блоками
    vertical_gap = Cm(0.5)  # Вертикальный зазор между строками (группами)

    current_top = top_start

    for ao_key, values_list in tops.items():
        current_left = left_group_box + group_box_width  # немного отступим вправо от группы
        match len(values_list):
            case 4:
                item_box_width = Cm(1.28)
                fontsize_main = Pt(12.4)
                fontsize_second = Pt(11)
            case 5:
                item_box_width = Cm(1.024)
                fontsize_main = Pt(9.4)
                fontsize_second = Pt(8)
            case _:
                item_box_width = Cm(1.7)
                fontsize_main = Pt(13.4)
                fontsize_second = Pt(12)
        for name, number in values_list:
            # Добавляем отдельный TextBox под элемент
            shape = slide.shapes.add_textbox(current_left, current_top, item_box_width, item_box_height)
            tf = shape.text_frame
            tf = shape.text_frame

            # Отключаем авто-изменение размера и ставим вертикальное выравнивание
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER

            run1 = p.add_run()
            run1.text = f"{name}\n"  # перевод строки
            run1.font.size = fontsize_main
            run1.font.name = "Golos UI Medium"
            run1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # белый цвет

            run2 = p.add_run()
            run2.text = str(number)
            run2.font.size = fontsize_second
            run2.font.name = "Golos UI Medium"
            run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            current_left += item_box_width + horizontal_gap
        current_top += group_box_height + vertical_gap

    return prs


def make_figure(new_slide, percents):
    # Цвета для фигур (до 4 категорий)
    colors_figure = [
        (159 << 16) + (135 << 8) + 97,  # #61879F
        (196 << 16) + (158 << 8) + 200,  # #C89EC4
        (197 << 16) + (127 << 8) + 117,  # #757FC5
        (146 << 16) + (214 << 8) + 253,  # #FDD692
    ]

    # Таблица топ-позиций в зависимости от количества элементов
    top_positions_dict = {
        4: [18.84, 21.36, 23.97, 23.97],
        3: [19.24, 22.73, 26.16],
        2: [19.78, 25.23],
        1: [22.73]
    }

    # Размеры и позиция (в поинтах, 1см = 28.35pt)
    width, height = 1.84 * 28.35, 1.07 * 28.35
    left = 1.41 * 28.35

    # Получаем нужные `top_positions`, если их нет — используем значения для 3 элементов
    top_positions = top_positions_dict.get(len(percents), top_positions_dict[3])

    # Создание фигур
    print(len(percents), percents)
    for i, percent in enumerate(percents):
        top = top_positions[min(i, len(top_positions) - 1)] * 28.35  # Выбираем top из списка
        pentagon_arrow = new_slide.Shapes.AddShape(51, left, top, width, height)

        # Настройки фигуры
        pentagon_arrow.TextFrame.TextRange.Text = f"{percent * 100}%".replace(".0", "")  # Округляем %
        if i > 3:
            i = 0
        pentagon_arrow.Fill.ForeColor.RGB = colors_figure[i]  # Цвет из списка
        pentagon_arrow.Fill.Transparency = 0.2  # Прозрачность 20%
        pentagon_arrow.Line.Visible = False  # Убираем обводку

        # Настройка текста
        text_range = pentagon_arrow.TextFrame.TextRange
        text_range.Font.Name = "Golos UI Medium"
        text_range.Font.Size = 16
        text_range.Font.Color.RGB = 0x000000  # Черный цвет текста


def make_yearly_slides(pptx_file_path, new_chart_datas, new_chart_datas_reasons, tables_regions):
    # Инициализация COM

    # df_regions = next(iter(tables_regions.values()))
    # event_key = next(key for key, value in tables_regions.items() if value is df_regions)
    # new_chart_data = df_to_chart_data(next(iter(new_chart_datas.values())))
    # new_chart_data_reasons = df_to_chart_data(next(iter(new_chart_datas_reasons.values())), True)
    comtypes.CoInitialize()

    try:
        ppt_path = os.path.abspath(pptx_file_path)

        # Открываем PowerPoint
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        ppt.Visible = 1  # 1 - показать PowerPoint, 0 - скрыть

        # Открываем презентацию
        presentation = ppt.Presentations.Open(ppt_path, WithWindow=False)
        page_id_rev = 0
        len_pages = len(new_chart_datas.items())
        for event, df in reversed(new_chart_datas.items()):
            event_key = event
            df_regions = tables_regions[event]
            new_chart_data = df_to_chart_data(df)
            new_chart_data_reasons = df_to_chart_data(new_chart_datas_reasons[event], True)
            # Индекс 5-го слайда (PowerPoint использует 1-based индексацию)
            slide_index = 5
            page_id_rev = page_id_rev + 1
            page_id = len_pages - page_id_rev + 5
            # Проверяем, что слайд существует
            if slide_index <= presentation.Slides.Count:
                slide_to_copy = presentation.Slides(slide_index)

                # Дублируем слайд
                new_slide = slide_to_copy.Duplicate()

                # Перемещаем дубликат после оригинального слайда
                new_slide.MoveTo(slide_index + 1)
                for shape_id, shape in enumerate(new_slide.Shapes):

                    if shape.HasChart and shape_id == 2:  # Проверяем, что объект — это диаграмма
                        chart = shape.Chart
                        chart_data = chart.ChartData

                        # Открываем данные диаграммы в Excel
                        chart_data.Workbook.Application.Visible = False
                        worksheet = chart_data.Workbook.Worksheets(1)

                        # Очищаем текущие данные
                        worksheet.Cells.ClearContents()

                        # Записываем новые данные (предполагается, что у диаграммы 1 серия)
                        worksheet.Cells(1, 1).FormulaR1C1 = "Категории"
                        worksheet.Cells(1, 2).FormulaR1C1 = "Значения"

                        sorted_data = sorted(new_chart_data, key=lambda x: x[1], reverse=True)  # Сортируем по значениям
                        top_3_categories = {category for category, _ in sorted_data[:3]}  # Берем первые 3 категории
                        sum_values = 0
                        for i, (category, value) in enumerate(new_chart_data, start=2):
                            worksheet.Cells(i, 1).FormulaR1C1 = category
                            worksheet.Cells(i, 2).FormulaR1C1 = value
                            sum_values = sum_values + value
                        # Закрываем таблицу данных
                        chart_data.Workbook.Close()
                        series = chart.SeriesCollection(1)  # Берем первую серию данных
                        for point_idx, point in enumerate(series.Points(), start=0):
                            category_name = new_chart_data[point_idx][0]  # Получаем название категории
                            if category_name in top_3_categories:
                                point.Format.Line.Visible = True  # Включаем контур
                                point.Format.Line.ForeColor.RGB = (0 << 16) + (
                                        0 << 8) + 255  # Красный цвет (RGB: 255, 0, 0)
                                point.Format.Line.Weight = 1.5
                    if shape.HasChart and shape_id == 1:
                        chart = shape.Chart
                        chart_data = chart.ChartData

                        # Открываем данные диаграммы в Excel
                        chart_data.Workbook.Application.Visible = False
                        worksheet = chart_data.Workbook.Worksheets(1)

                        # Очищаем текущие данные
                        worksheet.Cells.ClearContents()

                        # Записываем новые данные (предполагается, что у диаграммы 1 серия)
                        worksheet.Cells(1, 1).FormulaR1C1 = "Категории"
                        worksheet.Cells(1, 2).FormulaR1C1 = "Значения"

                        num_rows = len(new_chart_data_reasons)  # Количество заполненных строк
                        percent_values = []
                        for i, (category, value) in enumerate(new_chart_data_reasons, start=2):
                            worksheet.Cells(i, 1).FormulaR1C1 = category
                            worksheet.Cells(i, 2).FormulaR1C1 = value
                            percent_values.append(value)
                        if num_rows > 0:
                            last_row = num_rows + 1  # Так как данные начинаются со 2-й строки
                            sheet_name = worksheet.Name

                            data_range_str = f"{sheet_name}!A1:B{last_row}"
                            chart.SetSourceData(data_range_str)
                        make_figure(new_slide, percent_values)

                        # Закрываем таблицу данных
                        chart_data.Workbook.Close()

                    if shape.HasTable:  # Проверяем, является ли объект таблицей

                        table = shape.Table
                        row_count = len(df_regions)  # Количество строк данных
                        col_count = 3  # Количество колонок (название округа, процент, обращения на 1000 жителей)
                        # Вставляем данные
                        for i, (index, row) in enumerate(df_regions.iterrows(),
                                                         start=2):  # Строки начинаются со 2-й (первая - заголовок)
                            table.Cell(i, 1).Shape.TextFrame.TextRange.Text = row["Округ"]
                            table.Cell(i, 2).Shape.TextFrame.TextRange.Text = row["Процент от общего"]
                            table.Cell(i, 3).Shape.TextFrame.TextRange.Text = str(
                                row["Обращений на 1000 жителей"])  # Преобразуем в строку

                        print(f"✅ Данные успешно добавлены в таблицу на слайде {slide_index}")
                    if shape.HasTextFrame:  # Проверяем, является ли объект текстовым
                        text_range = shape.TextFrame.TextRange
                        if "all" in text_range.Text:  # Если в тексте есть "test"
                            text_range.Text = text_range.Text.replace("all", f"{sum_values}")  # Заменяем
                        if "page_id" in text_range.Text:  # Если в тексте есть "test"
                            text_range.Text = text_range.Text.replace("page_id", f"{page_id}")  # Заменяем
                    if shape.Type == 6:  # 6 - это тип GroupShape
                        for sub_shape in shape.GroupItems:  # Перебираем объекты внутри группы
                            if sub_shape.HasTextFrame:  # Проверяем, содержит ли объект текст
                                text_range = sub_shape.TextFrame.TextRange
                                if "test" in text_range.Text:  # Если в тексте есть "test"
                                    text_range.Text = text_range.Text.replace("test", f"{event_key}")  # Заменяем текст

                # Сохраняем изменения
            presentation.Save()

        # Закрываем презентацию и PowerPoint
        presentation.Close()
        ppt.Quit()

    finally:
        comtypes.CoUninitialize()

    return True  # Успешное выполнение

# 34 ЦАО 1
# 17 САО 2
# 14 СВАО 3
# 11 ВАО 4
# 46 ЮВАО 5
# 42 ЮАО 6
# 38 ЮЗАО 7
# 5 ЗАО 8
# 8 СЗАО 9
# 24 ЗелАО 10
# 30 ТиНАО 11


# 41 ЦАО
# 26 САО
# 23 СВАО
# 20 ВАО
# 53 ЮВАО
# 49 ЮАО
# 45 ЮЗАО
# 14 ЗАО
# 17 СЗАО
# 33 ЗелАО
# 37 ТиНАО
